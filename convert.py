import os
import tkinter as tk
from tkinter import filedialog, messagebox
from fpdf import FPDF
from PIL import Image
import pdfkit
import docx
import pandas as pd
from pptx import Presentation


FONT_PATH = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")
# Funzioni di conversione
def txt_to_pdf(input_path, output_path):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
    pdf.set_font("DejaVu", "", 12)

    with open(input_path, "r", encoding="utf-8") as file:
        for line in file:
            pdf.cell(200, 10, txt=line, ln=True)

    pdf.output(output_path)

def image_to_pdf(input_path, output_path):
    image = Image.open(input_path)
    image.convert("RGB").save(output_path)

def docx_to_pdf(input_path, output_path):
    doc = docx.Document(input_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
    pdf.set_font("DejaVu", "", 12)

    temp_dir = os.path.join(os.path.dirname(output_path), "temp_images")
    os.makedirs(temp_dir, exist_ok=True)

    y_position = 10 

    for para in doc.paragraphs:
        # Controllare lo stile del paragrafo
        if para.style.name.startswith("Heading"):
            pdf.set_font("DejaVu", "B", 14)  
            pdf.set_font("DejaVu", "", 12)

        
        line = ""
        for run in para.runs:
            if run.bold:
                pdf.set_font("DejaVu", "B", 12)  
            elif run.italic:
                pdf.set_font("DejaVu", "I", 12)  
            else:
                pdf.set_font("DejaVu", "", 12)

            line += run.text

        pdf.multi_cell(0, 10, line)  
        y_position += 10

    
    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            img_part = doc.part.rels[rel].target_part
            img_data = img_part.blob
            img_path = os.path.join(temp_dir, f"image_{rel}.png")

            with open(img_path, "wb") as img_file:
                img_file.write(img_data)

            
            if y_position + 50 > 270:
                pdf.add_page()
                y_position = 10

       
            pdf.image(img_path, x=10, y=y_position, w=100)
            y_position += 60  

    for table in doc.tables:
        pdf.add_page()
        pdf.set_font("DejaVu", "B", 12)
        pdf.cell(200, 10, "Tabella", ln=True)
        pdf.set_font("DejaVu", "", 12)

        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            pdf.cell(200, 10, " | ".join(row_text), ln=True)

    pdf.output(output_path)

    for img_file in os.listdir(temp_dir):
        os.remove(os.path.join(temp_dir, img_file))
    os.rmdir(temp_dir)


def excel_to_pdf(input_path, output_path):
    df = pd.read_excel(input_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
    pdf.set_font("DejaVu", "", 12)

    for _, row in df.iterrows():
        pdf.cell(200, 10, txt=str(row.values), ln=True)

    pdf.output(output_path)

def pptx_to_pdf(input_path, output_path):
    prs = Presentation(input_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
    pdf.set_font("DejaVu", "", 12)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.cell(200, 10, txt=shape.text, ln=True)

    pdf.output(output_path)

def convert_to_pdf():
    file_path = filedialog.askopenfilename(title="Seleziona un file")
    if not file_path:
        return
    
    file_type = os.path.splitext(file_path)[1].lower()
    output_path = os.path.splitext(file_path)[0] + ".pdf"

    try:
        if file_type == ".txt":
            txt_to_pdf(file_path, output_path)
        elif file_type in [".jpg", ".png", ".jpeg"]:
            image_to_pdf(file_path, output_path)
        elif file_type == ".docx":
            docx_to_pdf(file_path, output_path)
        elif file_type in [".xls", ".xlsx"]:
            excel_to_pdf(file_path, output_path)
        elif file_type == ".pptx":
            pptx_to_pdf(file_path, output_path)
        else:
            messagebox.showerror("Errore", "Formato non supportato")
            return

        messagebox.showinfo("Successo", f"Conversione completata! Salvato in:\n{output_path}")
    except Exception as e:
        messagebox.showerror("Errore", f"Errore durante la conversione:\n{e}")

# Creazione della GUI
root = tk.Tk()
root.title("Convertitore in PDF")
root.geometry("400x250")

label = tk.Label(root, text="Seleziona un file da convertire in PDF", font=("Arial", 12))
label.pack(pady=10)

button = tk.Button(root, text="Scegli un file", command=convert_to_pdf, font=("Arial", 12), bg="lightblue")
button.pack(pady=10)

exit_button = tk.Button(root, text="Esci", command=root.quit, font=("Arial", 12), bg="red", fg="white")
exit_button.pack(pady=10)

root.mainloop()
    
    # Il codice è abbastanza semplice. Abbiamo definito cinque funzioni di conversione per convertire file di testo, immagini, documenti Word, fogli di calcolo Excel e presentazioni PowerPoint in file PDF. 
    # La funzione  convert_to_pdf()  apre una finestra di dialogo per selezionare un file da convertire. In base all'estensione del file, viene chiamata la funzione di conversione corrispondente. 
    # Infine, abbiamo creato una GUI con una finestra di dialogo per selezionare un file e due pulsanti per avviare la conversione e chiudere l'applicazione. 
    # Eseguire il codice Python 
    # Per eseguire il codice Python, salvalo in un file chiamato  convert.py  e apri un terminale. Passa alla directory in cui hai salvato il file e digita il seguente comando per eseguire il programma: 
    # python convert.py 
    # Verrà visualizzata una finestra di dialogo per selezionare un file da convertire. Scegli un file e premi il pulsante "Scegli un file" per avviare la conversione. 
    # Una volta completata la conversione, verrà visualizzata una finestra di dialogo con un messaggio di successo e il percorso del file PDF convertito. 
    # Conclusioni 
    # In questo articolo, abbiamo visto come creare un convertitore di file in Python utilizzando la libreria  fpdf  per generare file PDF da file di testo, immagini, documenti Word, fogli di calcolo Excel e presentazioni PowerPoint. Abbiamo anche utilizzato la libreria  tkinter  per creare una GUI per l'applicazione. 
    # Questo è solo un esempio di come si possono creare applicazioni di conversione di file in Python. È possibile estendere il programma per supportare ulteriori formati di file o aggiungere funzionalità come la conversione batch di file.